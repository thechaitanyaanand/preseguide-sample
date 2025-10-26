import fitz  # PyMuPDF
from pptx import Presentation as PPTXPresentation
from typing import Dict, List
import os


class DocumentProcessor:
    """
    Process PDF and PowerPoint files to extract text and key points.
    """

    def process_document(self, file_path: str) -> Dict:
        """
        Main method to process a document file.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary with extracted text and key points
        """
        try:
            _, ext = os.path.splitext(file_path.lower())
            
            if ext == '.pdf':
                return self.process_pdf(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self.process_pptx(file_path)
            else:
                return {
                    'success': False,
                    'error': f'Unsupported file format: {ext}'
                }
                
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }

    def process_pdf(self, file_path: str) -> Dict:
        """Extract text from PDF."""
        try:
            doc = fitz.open(file_path)
            text = ""
            slides = []
            
            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text()
                text += page_text + "\n\n"
                
                slides.append({
                    'slide_number': page_num,
                    'content': page_text.strip(),
                    'word_count': len(page_text.split())
                })
            
            doc.close()
            
            # Extract key points (simple approach: look for bullet points or short lines)
            key_points = self.extract_key_points(text)
            
            return {
                'success': True,
                'full_text': text.strip(),
                'slides': slides,
                'key_points': key_points,
                'total_pages': len(slides),
                'total_words': len(text.split())
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'PDF processing error: {str(e)}'
            }

    def process_pptx(self, file_path: str) -> Dict:
        """Extract text from PowerPoint."""
        try:
            prs = PPTXPresentation(file_path)
            text = ""
            slides = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                text += slide_text + "\n\n"
                
                slides.append({
                    'slide_number': slide_num,
                    'content': slide_text.strip(),
                    'word_count': len(slide_text.split())
                })
            
            # Extract key points
            key_points = self.extract_key_points(text)
            
            return {
                'success': True,
                'full_text': text.strip(),
                'slides': slides,
                'key_points': key_points,
                'total_pages': len(slides),
                'total_words': len(text.split())
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'PowerPoint processing error: {str(e)}'
            }

    def extract_key_points(self, text: str) -> List[str]:
        """
        Extract key points from text.
        Simple approach: Look for short lines that might be headers or bullet points.
        """
        lines = text.split('\n')
        key_points = []
        
        for line in lines:
            line = line.strip()
            
            # Skip empty lines
            if not line:
                continue
            
            # Look for bullet points or short statements (5-50 words)
            words = line.split()
            if 5 <= len(words) <= 50:
                # Check if it starts with bullet indicators
                if any(line.startswith(char) for char in ['•', '-', '●', '○', '►']):
                    key_points.append(line)
                # Or if it's a short, declarative statement
                elif line[0].isupper() and not line.endswith(':'):
                    key_points.append(line)
        
        # Limit to top 20 key points
        return key_points[:20]

    def calculate_content_coverage(self, transcription: str, key_points: List[str]) -> Dict:
        """
        Calculate how well the transcription covers the key points.
        
        Returns:
            Dictionary with coverage score and missed points
        """
        if not key_points:
            return {
                'coverage_score': 100.0,
                'missed_points': []
            }
        
        transcription_lower = transcription.lower()
        covered_count = 0
        missed_points = []
        
        for point in key_points:
            # Check if key words from the point appear in transcription
            point_words = set(point.lower().split())
            
            # Remove common stop words
            stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for'}
            point_words = point_words - stop_words
            
            # If at least 50% of key words are in transcription, consider it covered
            if len(point_words) > 0:
                matches = sum(1 for word in point_words if word in transcription_lower)
                coverage_ratio = matches / len(point_words)
                
                if coverage_ratio >= 0.5:
                    covered_count += 1
                else:
                    missed_points.append(point)
        
        coverage_score = (covered_count / len(key_points)) * 100
        
        return {
            'coverage_score': round(coverage_score, 1),
            'missed_points': missed_points
        }
